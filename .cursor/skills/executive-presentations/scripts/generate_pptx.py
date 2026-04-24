#!/usr/bin/env python3
"""
generate_pptx.py — decks executivos com identidade de marca, narrativa (kicker/SCQA)
e layouts visuais mais fortes (hero, capítulo, cartões).

Uso:
    python generate_pptx.py deck.json output.pptx
    python generate_pptx.py deck.json out.pptx --fetch-images   # baixa image_url → temp

No JSON raiz, opcional:
  "brand": { "primary", "primary_dark", "secondary", "accent", "surface", "text_dark", "muted", "on_dark", "header_font", "body_font", "logo_path" }
  ou "brand_path": "brand.json"

Slides podem ter:
  "kicker": "Situação" | "Complicação" | "Dados" | "Conclusão"  (aparece no faixa superior)
  "image_path" | "image_url" (com --fetch-images)
  "layout": "title" | "chapter" | "scqa_story" | "hero_split" | "executive_summary" | ...
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import sys
import tempfile
import urllib.request
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

PALETTES: dict[str, dict[str, str]] = {
    "midnight_executive": {"primary": "1E2761", "secondary": "CADCFC", "accent": "F5A623", "surface": "F4F7FB", "text_dark": "1E2761", "muted": "6B7280"},
    "charcoal_minimal":   {"primary": "36454F", "secondary": "E8ECF0", "accent": "C45C26", "surface": "FAFAFA", "text_dark": "212121", "muted": "6B7280"},
    "ocean_gradient":     {"primary": "065A82", "secondary": "D4EAF4", "accent": "FF8C42", "surface": "F0F7FB", "text_dark": "0A2540", "muted": "5A6C7D"},
}

PROFILE_DEFAULTS: dict[str, dict[str, Any]] = {
    "corporate":  {"palette": "midnight_executive", "header_font": "Calibri Light", "body_font": "Calibri"},
    "startup":    {"palette": "ocean_gradient",     "header_font": "Calibri Light", "body_font": "Calibri"},
    "consulting": {"palette": "charcoal_minimal",   "header_font": "Calibri",       "body_font": "Calibri"},
}


def _hex_rgb(h: str) -> RGBColor:
    h = h.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _darken(c: RGBColor, factor: float = 0.52) -> RGBColor:
    def ch(x: int) -> int:
        return max(0, min(255, int(round(x * factor))))
    return RGBColor(ch(c[0]), ch(c[1]), ch(c[2]))


@dataclass
class Theme:
    primary: RGBColor
    primary_dark: RGBColor
    secondary: RGBColor
    accent: RGBColor
    surface: RGBColor
    text_on_light: RGBColor
    muted: RGBColor
    on_dark: RGBColor
    header_font: str
    body_font: str
    logo_path: str | None

    @classmethod
    def from_deck(cls, deck: dict) -> "Theme":
        brand: dict[str, Any] = dict(deck.get("brand") or {})
        bp = deck.get("brand_path")
        if bp and Path(bp).exists():
            brand.update(json.loads(Path(bp).read_text(encoding="utf-8")))
        profile = deck.get("profile", "corporate")
        prof = PROFILE_DEFAULTS.get(profile, PROFILE_DEFAULTS["corporate"])
        pal_name = deck.get("palette") or prof["palette"]
        pal = PALETTES.get(pal_name, PALETTES["midnight_executive"])

        if brand.get("primary"):
            primary = _hex_rgb(str(brand["primary"]))
            primary_dark = _hex_rgb(str(brand["primary_dark"])) if brand.get("primary_dark") else _darken(primary, 0.5)
            secondary = _hex_rgb(str(brand["secondary"])) if brand.get("secondary") else _hex_rgb(pal["secondary"])
            accent = _hex_rgb(str(brand["accent"])) if brand.get("accent") else _hex_rgb(pal["accent"])
            surface = _hex_rgb(str(brand["surface"])) if brand.get("surface") else _hex_rgb(pal["surface"])
            text_on_light = _hex_rgb(str(brand["text_dark"])) if brand.get("text_dark") else _hex_rgb(pal["text_dark"])
            muted = _hex_rgb(str(brand["muted"])) if brand.get("muted") else _hex_rgb(pal["muted"])
            on_dark = _hex_rgb(str(brand["on_dark"])) if brand.get("on_dark") else RGBColor(0xFF, 0xFF, 0xFF)
            hf = str(brand.get("header_font") or deck.get("header_font") or prof["header_font"])
            bf = str(brand.get("body_font") or deck.get("body_font") or prof["body_font"])
            logo = brand.get("logo_path") or deck.get("logo_path")
        else:
            primary = _hex_rgb(pal["primary"])
            primary_dark = _darken(primary, 0.48)
            secondary = _hex_rgb(pal["secondary"])
            accent = _hex_rgb(pal["accent"])
            surface = _hex_rgb(pal["surface"])
            text_on_light = _hex_rgb(pal["text_dark"])
            muted = _hex_rgb(pal["muted"])
            on_dark = RGBColor(0xFF, 0xFF, 0xFF)
            hf = deck.get("header_font") or prof["header_font"]
            bf = deck.get("body_font") or prof["body_font"]
            logo = deck.get("logo_path")

        return cls(
            primary=primary,
            primary_dark=primary_dark,
            secondary=secondary,
            accent=accent,
            surface=surface,
            text_on_light=text_on_light,
            muted=muted,
            on_dark=on_dark,
            header_font=hf,
            body_font=bf,
            logo_path=str(logo) if logo else None,
        )


# ------------------------------ helpers ------------------------------

def _set_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_rect(slide, left, top, w, h, fill_rgb: RGBColor, line=False) -> Any:
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    if not line:
        sh.line.fill.background()
    return sh


def _add_rounded(slide, left, top, w, h, fill_rgb: RGBColor) -> Any:
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    sh.line.color.rgb = fill_rgb
    sh.line.width = Pt(0)
    return sh


def _add_text(slide, left, top, width, height, text: str, *,
              font_name: str, size: int, bold: bool = False, color: RGBColor,
              align: PP_ALIGN = PP_ALIGN.LEFT, italic: bool = False,
              valign: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Emu(45720)
    tf.margin_right = Emu(45720)
    tf.margin_top = Emu(22860)
    tf.margin_bottom = Emu(22860)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_bullets(slide, left, top, width, height, bullets: list[str], *,
                 font_name: str, size: int, color: RGBColor, bullet_color: RGBColor) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(11)
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.name = font_name
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = txt
        r2.font.name = font_name
        r2.font.size = Pt(size)
        r2.font.color.rgb = color


def _source(slide, theme: Theme, text: str | None) -> None:
    if not text:
        return
    _add_text(
        slide, Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.35),
        f"Fonte: {text}",
        font_name=theme.body_font, size=8, color=theme.muted, italic=True,
    )


def _kicker_title_block(slide, theme: Theme, spec: dict, title_width: Any) -> None:
    kicker = spec.get("kicker") or spec.get("narrative_beat")
    if kicker:
        _add_text(
            slide, Inches(0.65), Inches(0.42), title_width, Inches(0.35),
            str(kicker).upper(),
            font_name=theme.body_font, size=10, bold=True, color=theme.accent,
        )
        _add_text(
            slide, Inches(0.65), Inches(0.78), title_width, Inches(1.35),
            spec.get("action_title", ""),
            font_name=theme.header_font, size=26, bold=True, color=theme.primary_dark,
        )
    else:
        _add_text(
            slide, Inches(0.65), Inches(0.48), title_width, Inches(1.55),
            spec.get("action_title", ""),
            font_name=theme.header_font, size=28, bold=True, color=theme.primary_dark,
        )


def _maybe_logo(slide, theme: Theme) -> None:
    if theme.logo_path and Path(theme.logo_path).exists():
        slide.shapes.add_picture(theme.logo_path, Inches(11.85), Inches(0.38), height=Inches(0.52))


def _resolve_image(spec: dict, fetch: bool) -> str | None:
    p = spec.get("image_path")
    if p and Path(p).exists():
        return str(Path(p).resolve())
    url = spec.get("image_url")
    if fetch and url:
        cache = Path(tempfile.gettempdir()) / "exec_pptx_imgs"
        cache.mkdir(parents=True, exist_ok=True)
        name = hashlib.sha256(url.encode()).hexdigest()[:16] + ".img"
        dest = cache / name
        if not dest.exists():
            urllib.request.urlretrieve(url, dest)  # noqa: S310
        return str(dest)
    return None


def _bento_visual(slide, left, top, w, h, theme: Theme) -> None:
    """Colagem geométrica quando não há foto — ainda assim há 'imagem' de marca."""
    _add_rect(slide, left, top, w, h, theme.secondary)
    pad = Inches(0.25)
    _add_rounded(slide, left + pad, top + pad, w * 0.42, h * 0.38, theme.primary)
    _add_rounded(slide, left + w * 0.48, top + pad * 2, w * 0.44, h * 0.28, theme.accent)
    _add_rounded(slide, left + pad, top + h * 0.48, w * 0.55, h * 0.42, _darken(theme.primary, 0.65))


def _place_photo_or_bento(slide, spec: dict, theme: Theme, left, top, w, h, fetch: bool) -> None:
    img = _resolve_image(spec, fetch)
    if img:
        slide.shapes.add_picture(img, left, top, width=w, height=h)
    else:
        _bento_visual(slide, left, top, w, h, theme)


# ------------------------------ layouts ------------------------------

def render_title(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    img_left = Inches(5.05)
    _place_photo_or_bento(slide, spec, theme, img_left, 0, SLIDE_W - img_left, SLIDE_H, fetch)
    panel = _add_rect(slide, 0, 0, Inches(5.35), SLIDE_H, theme.primary_dark)
    panel.line.fill.background()
    _add_rect(slide, Inches(5.25), 0, Inches(0.12), SLIDE_H, theme.accent)

    deck_title = spec.get("deck_tag") or spec.get("subtitle_prefix")
    if deck_title:
        _add_text(
            slide, Inches(0.65), Inches(0.55), Inches(4.5), Inches(0.4),
            str(deck_title).upper(),
            font_name=theme.body_font, size=9, bold=True, color=theme.accent,
        )
    _add_text(
        slide, Inches(0.65), Inches(1.15), Inches(4.5), Inches(2.2),
        spec.get("title", ""),
        font_name=theme.header_font, size=40, bold=True, color=theme.on_dark,
    )
    if spec.get("subtitle"):
        _add_text(
            slide, Inches(0.65), Inches(3.55), Inches(4.5), Inches(1.8),
            spec["subtitle"],
            font_name=theme.body_font, size=15, color=theme.secondary,
        )
    foot = []
    if spec.get("date"):
        foot.append(spec["date"])
    if spec.get("confidential"):
        foot.append("CONFIDENCIAL")
    if foot:
        _add_text(
            slide, Inches(0.65), Inches(6.95), Inches(4.5), Inches(0.35),
            " · ".join(foot),
            font_name=theme.body_font, size=9, color=theme.muted,
        )


def render_chapter(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.primary)
    num = spec.get("chapter_num", spec.get("chapter", "01"))
    _add_text(
        slide, Inches(0.55), Inches(0.45), Inches(3.5), Inches(1.6),
        str(num),
        font_name=theme.header_font, size=96, bold=True, color=_darken(theme.secondary, 0.75),
    )
    _add_text(
        slide, Inches(0.55), Inches(2.05), Inches(11.5), Inches(1.0),
        spec.get("chapter_title", spec.get("title", "")),
        font_name=theme.header_font, size=36, bold=True, color=theme.on_dark,
    )
    if spec.get("hook"):
        _add_text(
            slide, Inches(0.55), Inches(3.35), Inches(11.0), Inches(1.2),
            spec["hook"],
            font_name=theme.body_font, size=18, color=theme.secondary,
        )
    img = _resolve_image(spec, fetch)
    if img:
        slide.shapes.add_picture(img, Inches(9.2), Inches(1.0), width=Inches(3.9), height=Inches(5.2))
    else:
        _bento_visual(slide, Inches(9.2), Inches(1.0), Inches(3.9), Inches(5.2), theme)


def render_scqa_story(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    band = _add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), theme.primary_dark)
    band.line.fill.background()
    _add_text(
        slide, Inches(0.65), Inches(0.35), Inches(11.5), Inches(0.55),
        (spec.get("kicker") or "HISTÓRIA DOS DADOS").upper(),
        font_name=theme.body_font, size=10, bold=True, color=theme.accent,
    )
    _add_text(
        slide, Inches(0.65), Inches(0.62), Inches(11.5), Inches(0.65),
        spec.get("action_title", "SCQA"),
        font_name=theme.header_font, size=22, bold=True, color=theme.on_dark,
    )
    cards = [
        ("S", spec.get("situation", "")),
        ("C", spec.get("complication", "")),
        ("Q", spec.get("question", "")),
        ("A", spec.get("answer", "")),
    ]
    gx, gy = Inches(0.55), Inches(1.45)
    cw, ch = Inches(6.0), Inches(2.65)
    gap = Inches(0.35)
    positions = [(gx, gy), (gx + cw + gap, gy), (gx, gy + ch + gap * 0.5), (gx + cw + gap, gy + ch + gap * 0.5)]
    for (letter, body), (left, top) in zip(cards, positions):
        card = _add_rounded(slide, left, top, cw, ch, theme.surface)
        card.line.color.rgb = theme.secondary
        card.line.width = Pt(1)
        _add_rect(slide, left + Inches(0.2), top + Inches(0.22), Inches(0.5), Inches(0.5), theme.accent)
        _add_text(slide, left + Inches(0.28), top + Inches(0.24), Inches(0.4), Inches(0.45), letter,
                  font_name=theme.header_font, size=18, bold=True, color=theme.primary_dark)
        _add_text(slide, left + Inches(0.85), top + Inches(0.25), cw - Inches(1.05), ch - Inches(0.45), body,
                  font_name=theme.body_font, size=13, color=theme.text_on_light)
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_hero_split(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    left_w = Inches(6.15)
    _add_rect(slide, 0, 0, left_w, SLIDE_H, theme.surface)
    _kicker_title_block(slide, theme, spec, left_w - Inches(0.35))
    if spec.get("hero_stat"):
        _add_text(
            slide, Inches(0.6), Inches(2.55), left_w - Inches(0.9), Inches(1.4),
            str(spec["hero_stat"]),
            font_name=theme.header_font, size=64, bold=True, color=theme.accent,
        )
    if spec.get("hero_label"):
        _add_text(
            slide, Inches(0.6), Inches(3.85), left_w - Inches(0.9), Inches(0.55),
            str(spec["hero_label"]),
            font_name=theme.body_font, size=12, color=theme.muted,
        )
    if spec.get("bullets"):
        _add_bullets(
            slide, Inches(0.6), Inches(4.45), left_w - Inches(0.85), Inches(2.35),
            spec["bullets"][:3],
            font_name=theme.body_font, size=14, color=theme.text_on_light, bullet_color=theme.primary,
        )
    _place_photo_or_bento(slide, spec, theme, left_w, 0, SLIDE_W - left_w, SLIDE_H, fetch)
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_executive_summary(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    _kicker_title_block(slide, theme, spec, Inches(7.0))
    bullets = spec.get("bullets", [])
    _add_bullets(
        slide, Inches(0.65), Inches(2.35), Inches(7.0), Inches(4.35),
        bullets,
        font_name=theme.body_font, size=15, color=theme.text_on_light, bullet_color=theme.primary,
    )
    _place_photo_or_bento(slide, spec, theme, Inches(8.05), Inches(1.15), Inches(5.0), Inches(5.65), fetch)
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_kpi_callout(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    _add_rect(slide, 0, Inches(1.15), SLIDE_W, SLIDE_H - Inches(1.15), theme.secondary)
    band = _add_rect(slide, 0, 0, SLIDE_W, Inches(1.18), theme.primary_dark)
    band.line.fill.background()
    _kicker_title_block(slide, theme, spec, Inches(12.0))
    kpis = spec.get("kpis", [])[:3]
    col_w = Inches(11.8 / max(len(kpis), 1))
    top_val = Inches(2.35)
    for i, k in enumerate(kpis):
        left = Inches(0.75) + Emu(int(col_w) * i)
        card = _add_rounded(slide, left + Inches(0.08), Inches(2.05), col_w - Inches(0.2), Inches(4.35), theme.surface)
        card.line.color.rgb = theme.secondary
        card.line.width = Pt(1)
        color_val = theme.accent if k.get("highlight") else theme.primary
        _add_text(
            slide, left, top_val, col_w, Inches(1.55),
            str(k.get("value", "")),
            font_name=theme.header_font, size=56, bold=True, color=color_val,
            align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, left, Inches(4.15), col_w, Inches(1.55),
            str(k.get("label", "")),
            font_name=theme.body_font, size=11, color=theme.muted,
            align=PP_ALIGN.CENTER,
        )
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_two_column(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.22), theme.primary_dark)
    _kicker_title_block(slide, theme, spec, Inches(12.0))
    _add_bullets(
        slide, Inches(0.65), Inches(1.55), Inches(5.85), Inches(5.1),
        spec.get("left_bullets", []),
        font_name=theme.body_font, size=14, color=theme.text_on_light, bullet_color=theme.primary,
    )
    right = spec.get("right_visual") or {}
    if right.get("type") == "bar":
        _bar_chart(slide, Inches(6.85), Inches(1.55), Inches(6.35), Inches(4.85), right, theme)
    elif spec.get("right_text"):
        box = _add_rounded(slide, Inches(6.75), Inches(1.45), Inches(6.45), Inches(5.2), theme.surface)
        box.line.color.rgb = theme.secondary
        box.line.width = Pt(1)
        _add_text(
            slide, Inches(6.95), Inches(1.65), Inches(6.1), Inches(4.85),
            spec["right_text"],
            font_name=theme.body_font, size=13, color=theme.text_on_light,
        )
    else:
        _place_photo_or_bento(slide, spec, theme, Inches(6.75), Inches(1.45), Inches(6.45), Inches(5.2), fetch)
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_chart(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.22), theme.primary_dark)
    _kicker_title_block(slide, theme, spec, Inches(12.0))
    chart = spec.get("chart", {})
    ctype = chart.get("type", "bar")
    if ctype in ("bar", "column"):
        _column_chart(slide, Inches(0.65), Inches(1.45), Inches(8.35), Inches(5.25), chart, theme)
    elif ctype == "line":
        _line_chart(slide, Inches(0.65), Inches(1.45), Inches(8.35), Inches(5.25), chart, theme)
    elif ctype == "combo":
        _column_chart(
            slide, Inches(0.65), Inches(1.45), Inches(8.35), Inches(5.25),
            {"x": chart.get("x", []), "series": [chart.get("bars", {})]}, theme,
        )
    insight_spec = {**spec, "image_path": spec.get("insight_image_path") or spec.get("image_path")}
    _place_photo_or_bento(slide, insight_spec, theme, Inches(9.15), Inches(1.45), Inches(3.95), Inches(3.1), fetch)
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_comparison(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.22), theme.primary_dark)
    _kicker_title_block(slide, theme, spec, Inches(12.0))
    cols = spec.get("columns", [])
    rows = spec.get("rows", [])
    n_cols = 1 + len(cols)
    n_rows = 1 + len(rows)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.65), Inches(1.45), Inches(12.05), Inches(4.95))
    table = shape.table
    table.cell(0, 0).text = ""
    for j, c in enumerate(cols, start=1):
        name = c.get("name", "")
        if c.get("recommended"):
            name = f"★ {name}"
        cell = table.cell(0, j)
        cell.text = name
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.name = theme.header_font
                run.font.size = Pt(12)
                run.font.color.rgb = theme.on_dark if c.get("recommended") else theme.primary_dark
        try:
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.primary if c.get("recommended") else theme.secondary
        except Exception:
            pass
    for i, r in enumerate(rows, start=1):
        cell = table.cell(i, 0)
        cell.text = r.get("dimension", "")
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.name = theme.body_font
                run.font.size = Pt(11)
                run.font.color.rgb = theme.text_on_light
        vals = r.get("values", [])
        for j, v in enumerate(vals, start=1):
            c2 = table.cell(i, j)
            c2.text = str(v)
            for p in c2.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = theme.body_font
                    run.font.size = Pt(11)
                    run.font.color.rgb = theme.text_on_light
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_timeline(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.22), theme.primary_dark)
    _kicker_title_block(slide, theme, spec, Inches(12.0))
    ms = spec.get("milestones", [])
    if not ms:
        return
    track_top = Inches(4.05)
    left0 = Inches(0.95)
    track_w = Inches(11.5)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left0, track_top, track_w, Emu(38100))
    line.fill.solid()
    line.fill.fore_color.rgb = theme.muted
    line.line.fill.background()
    step = int(track_w) // max(len(ms) - 1, 1)
    for i, m in enumerate(ms):
        cx = int(left0) + step * i
        dot_size = Inches(0.34)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            cx - int(dot_size) // 2,
            int(track_top) - int(dot_size) // 2 + Emu(19050),
            dot_size, dot_size,
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = theme.accent if m.get("highlight") else theme.primary
        dot.line.fill.background()
        _add_text(
            slide, Emu(cx - Inches(1.15) // 2), Inches(2.35), Inches(1.15), Inches(0.45),
            m.get("when", ""),
            font_name=theme.body_font, size=10, color=theme.muted, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, Emu(cx - Inches(1.55) // 2), Inches(4.55), Inches(1.55), Inches(0.95),
            m.get("label", ""),
            font_name=theme.body_font, size=11, bold=True, color=theme.text_on_light, align=PP_ALIGN.CENTER,
        )
    _place_photo_or_bento(slide, spec, theme, Inches(9.0), Inches(1.5), Inches(4.0), Inches(2.8), fetch)
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_image_right(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.surface)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.22), theme.primary_dark)
    _kicker_title_block(slide, theme, spec, Inches(6.9))
    _add_bullets(
        slide, Inches(0.65), Inches(1.55), Inches(6.35), Inches(5.1),
        spec.get("bullets", []),
        font_name=theme.body_font, size=14, color=theme.text_on_light, bullet_color=theme.primary,
    )
    _place_photo_or_bento(slide, spec, theme, Inches(7.15), Inches(1.35), Inches(5.95), Inches(5.55), fetch)
    _maybe_logo(slide, theme)
    _source(slide, theme, spec.get("source"))


def render_closing(slide, spec: dict, theme: Theme, fetch: bool) -> None:
    _set_bg(slide, theme.primary_dark)
    _add_rect(slide, Inches(7.4), 0, Inches(5.93), SLIDE_H, theme.primary)
    _place_photo_or_bento(slide, spec, theme, Inches(7.4), 0, Inches(5.93), SLIDE_H, fetch)
    _add_text(
        slide, Inches(0.65), Inches(0.55), Inches(6.5), Inches(1.6),
        spec.get("action_title", ""),
        font_name=theme.header_font, size=30, bold=True, color=theme.on_dark,
    )
    _add_text(
        slide, Inches(0.65), Inches(2.35), Inches(6.2), Inches(0.45),
        "Próximos passos",
        font_name=theme.body_font, size=12, bold=True, color=theme.accent,
    )
    steps = spec.get("next_steps", [])
    for i, s in enumerate(steps):
        row_top = Inches(2.85) + Inches(0.52 * i)
        _add_text(slide, Inches(0.65), row_top, Inches(0.45), Inches(0.45), f"{i + 1}.",
                  font_name=theme.body_font, size=13, bold=True, color=theme.secondary)
        _add_text(slide, Inches(1.1), row_top, Inches(5.8), Inches(0.45), s.get("action", ""),
                  font_name=theme.body_font, size=13, color=theme.on_dark)
        _add_text(slide, Inches(6.5), row_top, Inches(2.5), Inches(0.45), s.get("owner", ""),
                  font_name=theme.body_font, size=12, color=theme.secondary, italic=True)
    if spec.get("next_review"):
        _add_text(
            slide, Inches(0.65), Inches(6.95), Inches(6.5), Inches(0.35),
            f"Próxima revisão: {spec['next_review']}",
            font_name=theme.body_font, size=9, color=theme.muted, italic=True,
        )


# ------------------------------ charts ------------------------------

def _column_chart(slide, left, top, width, height, chart_spec: dict, theme: Theme) -> None:
    data = CategoryChartData()
    cats = chart_spec.get("x") or chart_spec.get("categories") or []
    data.categories = cats
    series = chart_spec.get("series")
    if series is None:
        bars = chart_spec.get("bars") or chart_spec.get("values")
        if isinstance(bars, dict):
            data.add_series(bars.get("label", "Série 1"), bars.get("values", []))
        elif isinstance(bars, list):
            data.add_series("Série 1", bars)
    else:
        for s in series:
            data.add_series(s.get("label", "Série"), s.get("values", []))
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data
    ).chart
    _style_chart(graphic, theme)
    try:
        plot = graphic.plots[0]
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = theme.accent if i % 2 else theme.primary
    except Exception:
        pass


def _bar_chart(slide, left, top, width, height, spec: dict, theme: Theme) -> None:
    data = CategoryChartData()
    rows = spec.get("data", [])
    data.categories = [r.get("label", "") for r in rows]
    data.add_series(spec.get("unit", "Valor"), [r.get("value", 0) for r in rows])
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, data
    ).chart
    _style_chart(graphic, theme)
    plot = graphic.plots[0]
    for i, r in enumerate(rows):
        try:
            pt = plot.series[0].points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = theme.accent if r.get("highlight") else theme.primary
        except Exception:
            pass


def _line_chart(slide, left, top, width, height, chart_spec: dict, theme: Theme) -> None:
    data = CategoryChartData()
    data.categories = chart_spec.get("x", [])
    for s in chart_spec.get("series", []):
        data.add_series(s.get("label", "Série"), s.get("values", []))
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, data
    ).chart
    _style_chart(graphic, theme)


def _style_chart(chart, theme: Theme) -> None:
    chart.has_title = False
    try:
        chart.plots[0].has_data_labels = True
        dl = chart.plots[0].data_labels
        dl.font.size = Pt(9)
        dl.font.name = theme.body_font
        dl.font.color.rgb = theme.text_on_light
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass


# ------------------------------ dispatcher ------------------------------

def _dispatch(layout: str):
    return {
        "title": render_title,
        "chapter": render_chapter,
        "scqa_story": render_scqa_story,
        "hero_split": render_hero_split,
        "executive_summary": render_executive_summary,
        "kpi_callout": render_kpi_callout,
        "two_column": render_two_column,
        "chart": render_chart,
        "comparison": render_comparison,
        "timeline": render_timeline,
        "image_right": render_image_right,
        "closing": render_closing,
    }.get(layout, render_executive_summary)


def _expand_media_paths(deck: dict, deck_json: Path) -> None:
    """Resolve image_path / logo_path relativos à pasta do deck.json."""
    root = deck_json.parent

    def fix(p: str | None) -> str | None:
        if not p:
            return None
        path = Path(p)
        if not path.is_absolute():
            path = (root / path).resolve()
        return str(path) if path.exists() else p

    b = deck.get("brand") or {}
    if b.get("logo_path"):
        b["logo_path"] = fix(str(b["logo_path"]))
    if deck.get("logo_path"):
        deck["logo_path"] = fix(str(deck["logo_path"]))
    for spec in deck.get("slides", []):
        if spec.get("image_path"):
            spec["image_path"] = fix(str(spec["image_path"]))
        if spec.get("insight_image_path"):
            spec["insight_image_path"] = fix(str(spec["insight_image_path"]))


def build(deck: dict, output: Path, fetch_images: bool, deck_json: Path | None = None) -> None:
    if deck_json:
        _expand_media_paths(deck, deck_json)
    theme = Theme.from_deck(deck)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    for spec in deck.get("slides", []):
        layout = spec.get("layout", "executive_summary")
        renderer = _dispatch(layout)
        slide = prs.slides.add_slide(blank_layout)
        renderer(slide, spec, theme, fetch_images)

    prs.save(str(output))
    print(f"[ok] gerado: {output}  ({len(deck.get('slides', []))} slides)")


def main() -> int:
    parser = argparse.ArgumentParser(description="Gera deck executivo .pptx a partir de JSON.")
    parser.add_argument("deck_json", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument("--fetch-images", action="store_true", help="Baixa image_url da internet para slides")
    args = parser.parse_args()

    if not args.deck_json.exists():
        print(f"[erro] {args.deck_json} não encontrado", file=sys.stderr)
        return 1
    try:
        deck = json.loads(args.deck_json.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        print(f"[erro] JSON inválido: {e}", file=sys.stderr)
        return 1

    build(
        deck,
        args.output_pptx,
        args.fetch_images or os.environ.get("EXEC_PPTX_FETCH_IMAGES") == "1",
        deck_json=args.deck_json,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
